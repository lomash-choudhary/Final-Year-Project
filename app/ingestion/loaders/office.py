"""
DOCX / PPTX loader.

Uses python-docx and python-pptx directly rather than `unstructured`. Both are
tiny, pure-Python and have no system dependencies, whereas `unstructured` pulls a
large tree and frequently breaks on fresh installs. PPTX maps one slide to one
`Page`, which gives slide-accurate citations for free.
"""

from __future__ import annotations

from pathlib import Path

import logfire

from app.ingestion.loaders.base import ExtractionFailed, LoadedDocument, Page


def _load_docx(path: Path) -> list[Page]:
    from docx import Document as DocxDocument

    document = DocxDocument(str(path))
    parts = [p.text for p in document.paragraphs if p.text.strip()]

    # Tables carry most of the numbers in veterinary/epidemiology reports.
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    text = "\n\n".join(parts)
    return [Page(number=1, text=text, extractor="python-docx")] if text.strip() else []


def _load_pptx(path: Path) -> list[Page]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    pages: list[Page] = []

    for idx, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        text = "\n".join(parts)
        if text.strip():
            pages.append(Page(number=idx, text=text, extractor="python-pptx"))

    return pages


def load_office(path: Path) -> LoadedDocument:
    ext = path.suffix.lower().lstrip(".")

    with logfire.span("Office extraction", filename=path.name, format=ext):
        try:
            pages = _load_docx(path) if ext == "docx" else _load_pptx(path)
        except ImportError as exc:
            raise ExtractionFailed(
                f"Missing parser for .{ext}: {exc}. Install python-docx / python-pptx."
            ) from exc
        except Exception as exc:
            raise ExtractionFailed(f"Could not parse {path.name}: {exc}") from exc

        if not pages:
            raise ExtractionFailed(f"{path.name} produced no readable text.")

        doc = LoadedDocument(filename=path.name, path=str(path), loader=f"office-{ext}", pages=pages)
        logfire.info("Office document extracted", filename=path.name, pages=doc.page_count, chars=doc.char_count)
        return doc
